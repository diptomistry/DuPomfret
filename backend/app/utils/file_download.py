"""File download and text extraction utilities.

Extended to support code files (.py, .js, .java, .ts, etc.), Jupyter notebooks (.ipynb),
and ZIP archives containing code or notebooks.
"""
import httpx
import io
import json
import zipfile
import os
from typing import Optional, List
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation


async def download_file(url: str) -> bytes:
    """
    Download file from URL.
    """
    async with httpx.AsyncClient() as client:
        response = await client.get(url, timeout=30.0)
        response.raise_for_status()
        return response.content


def extract_text_from_pdf(file_content: bytes) -> str:
    pdf_file = io.BytesIO(file_content)
    reader = PdfReader(pdf_file)
    text_parts: List[str] = []
    for page in reader.pages:
        text_parts.append(page.extract_text() or "")
    return "\n".join(text_parts)


def extract_text_from_docx(file_content: bytes) -> str:
    docx_file = io.BytesIO(file_content)
    doc = Document(docx_file)
    text_parts: List[str] = []
    for paragraph in doc.paragraphs:
        text_parts.append(paragraph.text)
    return "\n".join(text_parts)


def extract_text_from_txt(file_content: bytes) -> str:
    try:
        return file_content.decode("utf-8")
    except UnicodeDecodeError:
        return file_content.decode("latin-1", errors="ignore")


def extract_text_from_ipynb(file_content: bytes) -> str:
    """
    Extract code and markdown from a Jupyter notebook (.ipynb).
    """
    try:
        nb = json.loads(file_content.decode("utf-8"))
    except Exception:
        # Fallback to tolerant decoding
        nb = json.loads(file_content.decode("latin-1", errors="ignore"))

    parts: List[str] = []
    for cell in nb.get("cells", []):
        ctype = cell.get("cell_type")
        source = "".join(cell.get("source", []) or [])
        if not source:
            continue
        if ctype == "markdown":
            parts.append(source)
        elif ctype == "code":
            # keep code blocks as-is
            parts.append(source)
    return "\n\n".join(parts)


def extract_text_from_pptx(file_content: bytes) -> str:
    """
    Extract text from a .pptx PowerPoint file using python-pptx.
    """
    buf = io.BytesIO(file_content)
    try:
        prs = Presentation(buf)
    except Exception:
        return ""
    parts: List[str] = []
    for slide in prs.slides:
        slide_text_parts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text or ""
                if text.strip():
                    slide_text_parts.append(text)
        if slide_text_parts:
            parts.append("\n".join(slide_text_parts))
    return "\n\n".join(parts)


def extract_pptx_slides(file_content: bytes) -> List[str]:
    """
    Return a list of slide texts for a .pptx file.
    """
    buf = io.BytesIO(file_content)
    try:
        prs = Presentation(buf)
    except Exception:
        return []
    slides: List[str] = []
    for slide in prs.slides:
        slide_text_parts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text or ""
                if text.strip():
                    slide_text_parts.append(text)
        slides.append("\n".join(slide_text_parts).strip())
    return slides


def extract_text_from_zip(file_content: bytes) -> str:
    """
    Iterate files inside a ZIP archive and concatenate supported file contents.
    """
    allowed = {
        ".py",
        ".js",
        ".ts",
        ".java",
        ".c",
        ".cpp",
        ".go",
        ".rs",
        ".html",
        ".css",
        ".md",
        ".json",
        ".yaml",
        ".yml",
        ".sh",
        ".ps1",
        ".txt",
        ".ipynb",
        ".docx",
        ".pptx",
        ".ppt",
        ".pdf",
    }
    zbuf = io.BytesIO(file_content)
    parts: List[str] = []
    try:
        with zipfile.ZipFile(zbuf) as z:
            for name in z.namelist():
                _, ext = os.path.splitext(name.lower())
                if ext not in allowed:
                    continue
                try:
                    data = z.read(name)
                except Exception:
                    continue
                if ext == ".pdf":
                    parts.append(extract_text_from_pdf(data))
                elif ext == ".docx":
                    parts.append(extract_text_from_docx(data))
                elif ext == ".pptx":
                    parts.append(extract_text_from_pptx(data))
                elif ext == ".ipynb":
                    parts.append(extract_text_from_ipynb(data))
                else:
                    parts.append(extract_text_from_txt(data))
    except zipfile.BadZipFile:
        return ""
    return "\n\n".join([p for p in parts if p])


async def extract_text_from_file(url: str) -> str:
    """
    Download file and extract text based on file extension.

    Supported formats:
      - Documents: PDF, DOCX, TXT
      - Code files: .py, .js, .ts, .java, .c, .cpp, .go, .rs, .html, .css, .md, .json, .yaml
      - Jupyter notebooks: .ipynb
      - Zip archives containing the above
    """
    file_content = await download_file(url)

    # strip query params/fragments
    path = url.split("?", 1)[0].split("#", 1)[0].lower()
    _, ext = os.path.splitext(path)

    if ext == ".pdf":
        return extract_text_from_pdf(file_content)
    elif ext in (".docx", ".doc"):
        return extract_text_from_docx(file_content)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_content)
    elif ext == ".ppt":
        # .ppt (legacy binary PowerPoint) is not directly supported. Recommend converting to .pptx.
        raise ValueError("Unsupported legacy .ppt format. Please convert to .pptx and try again.")
    elif ext == ".txt":
        return extract_text_from_txt(file_content)
    elif ext == ".ipynb":
        return extract_text_from_ipynb(file_content)
    elif ext == ".zip":
        return extract_text_from_zip(file_content)
    else:
        # Treat many common source-code file types as text
        code_exts = {
            ".py",
            ".js",
            ".ts",
            ".java",
            ".c",
            ".cpp",
            ".go",
            ".rs",
            ".html",
            ".css",
            ".md",
            ".json",
            ".yaml",
            ".yml",
            ".sh",
            ".ps1",
        }
        if ext in code_exts or ext == "":
            # If URL has no extension, still attempt to decode as text
            return extract_text_from_txt(file_content)
        raise ValueError(
            "Unsupported file format. Supported: PDF, DOCX, PPTX, TXT, code files, .ipynb, .zip"
        )
